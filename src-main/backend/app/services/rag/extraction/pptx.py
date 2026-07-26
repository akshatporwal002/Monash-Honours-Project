"""Slide-preserving PPTX text extraction."""

from __future__ import annotations

from typing import BinaryIO

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.services.rag.contracts import ExtractedBlock, ExtractedDocument
from app.services.rag.errors import InvalidDocumentError
from app.services.rag.extraction.base import document_from_blocks, validate_openxml_archive


class PptxDocumentExtractor:
    supported_mime_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )

    def extract(self, source: BinaryIO) -> ExtractedDocument:
        validate_openxml_archive(source, "ppt/")
        try:
            presentation = Presentation(source)
            blocks: list[ExtractedBlock] = []
            title = presentation.core_properties.title
            for slide_number, slide in enumerate(presentation.slides, start=1):
                heading: str | None = None
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text.strip()
                    if not text:
                        continue
                    is_title = bool(
                        shape.is_placeholder
                        and shape.placeholder_format.type
                        in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
                    )
                    if is_title:
                        heading = text
                    blocks.append(
                        ExtractedBlock(
                            ordinal=len(blocks),
                            text=text,
                            heading=heading,
                            location_label=f"Slide {slide_number}",
                            block_type="heading" if is_title else "paragraph",
                        )
                    )
            return document_from_blocks(blocks, title)
        except (KeyError, OSError, ValueError) as error:
            raise InvalidDocumentError() from error
