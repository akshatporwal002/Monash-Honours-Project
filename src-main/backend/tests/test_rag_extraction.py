import io
import zipfile

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from app.services.rag.errors import (
    EncryptedDocumentError,
    InvalidDocumentError,
    NoExtractableTextError,
    UnsupportedMaterialTypeError,
)
from app.services.rag.extraction.docx import DocxDocumentExtractor
from app.services.rag.extraction.pdf import PdfDocumentExtractor
from app.services.rag.extraction.pptx import PptxDocumentExtractor
from app.services.rag.extraction.registry import ExtractorRegistry


def _pdf_with_text(text: str) -> bytes:
    content = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, value in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{number} 0 obj\n".encode() + value + b"\nendobj\n")
    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode())
    output.extend(b"".join(f"{offset:010d} 00000 n \n".encode() for offset in offsets[1:]))
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF".encode()
    )
    return bytes(output)


def _docx_bytes() -> bytes:
    document = Document()
    document.core_properties.title = "Quantum notes"
    document.add_heading("Superposition", level=1)
    document.add_paragraph("A Hadamard gate creates superposition.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Gate"
    table.cell(0, 1).text = "H"
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Measurement"
    slide.placeholders[1].text = "Measurement gives a classical outcome."
    extra = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
    extra.text_frame.text = "Born rule"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_pdf_extractor_preserves_page_text_and_rejects_empty_or_encrypted_documents() -> None:
    extractor = PdfDocumentExtractor()
    extracted = extractor.extract(io.BytesIO(_pdf_with_text("Quantum state")))

    assert extracted.blocks[0].text == "Quantum state"
    assert extracted.blocks[0].location_label == "Page 1"

    blank = PdfWriter()
    blank.add_blank_page(width=72, height=72)
    blank_bytes = io.BytesIO()
    blank.write(blank_bytes)
    with pytest.raises(NoExtractableTextError):
        extractor.extract(io.BytesIO(blank_bytes.getvalue()))

    encrypted = PdfWriter()
    encrypted.add_blank_page(width=72, height=72)
    encrypted.encrypt("secret")
    encrypted_bytes = io.BytesIO()
    encrypted.write(encrypted_bytes)
    with pytest.raises(EncryptedDocumentError):
        extractor.extract(io.BytesIO(encrypted_bytes.getvalue()))


def test_docx_extractor_preserves_order_headings_and_tables() -> None:
    extracted = DocxDocumentExtractor().extract(io.BytesIO(_docx_bytes()))

    assert extracted.title == "Quantum notes"
    assert [block.block_type for block in extracted.blocks] == ["heading", "paragraph", "table"]
    assert extracted.blocks[1].heading == "Superposition"
    assert extracted.blocks[2].text == "Gate | H"


def test_pptx_extractor_preserves_slide_text_and_title_heading() -> None:
    extracted = PptxDocumentExtractor().extract(io.BytesIO(_pptx_bytes()))

    assert extracted.blocks[0].location_label == "Slide 1"
    assert extracted.blocks[0].heading == "Measurement"
    assert "Born rule" in [block.text for block in extracted.blocks]


def test_openxml_extractors_reject_corrupt_or_archive_bomb_inputs() -> None:
    with pytest.raises(InvalidDocumentError):
        DocxDocumentExtractor().extract(io.BytesIO(b"not a zip"))

    bomb = io.BytesIO()
    with zipfile.ZipFile(bomb, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "x")
        archive.writestr("word/document.xml", "a" * (1024 * 1024))
    with pytest.raises(InvalidDocumentError):
        DocxDocumentExtractor().extract(io.BytesIO(bomb.getvalue()))


def test_registry_selects_by_mime_type_and_rejects_unknown_types() -> None:
    registry = ExtractorRegistry(
        [PdfDocumentExtractor(), DocxDocumentExtractor(), PptxDocumentExtractor()]
    )

    assert isinstance(registry.get("application/pdf"), PdfDocumentExtractor)
    with pytest.raises(UnsupportedMaterialTypeError):
        registry.get("text/plain")
